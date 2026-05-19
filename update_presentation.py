"""
update_presentation.py — Update מצגת slide 18 with billing data.

Usage:
  python update_presentation.py --src source.pptx --dest updated.pptx
  python update_presentation.py   # uses defaults below
"""

import argparse
import sys
import copy
from pathlib import Path
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

_DEFAULT_SRC  = Path("presentations") / "source.pptx"
_DEFAULT_DEST = Path("presentations") / "updated.pptx"

_parser = argparse.ArgumentParser(description="Update PPTX slide 18 with billing data")
_parser.add_argument("--src",  type=Path, default=_DEFAULT_SRC,  help="Source .pptx path")
_parser.add_argument("--dest", type=Path, default=_DEFAULT_DEST, help="Output .pptx path")
_args, _ = _parser.parse_known_args()

SRC  = str(_args.src)
DEST = str(_args.dest)

# ── March 2026 calculated values ────────────────────────────────────────────
# Basis: actual rate 34.62₪/h, avg hours: 168.20 @ 100% + 40.01 @ 125% + 30.95 @ 150%
# Gross = 5,821.77 + 1,731.23 + 1,607.62 = 9,160.62
# Standard 240h used for cost-per-employee model; hourly = total/240

TOP_ROW = {
    # col: new_text
    1:  " ₪ 9,160.62 ",   # שכר ברוטו
    2:  " ₪   687.05 ",   # פנסיה 7.5%
    3:  " ₪   109.93 ",   # ביטוח לאומי מעסיק
    4:  " ₪   288.56 ",   # עלות חופשים 3.15%
    5:  " ₪   618.32 ",   # אגרה שנתית (unchanged)
    6:  " ₪ 1,047.06 ",   # אגרה דמי היתר (unchanged)
    7:  " ₪    36.72 ",   # קרן עידוד (actual from data)
    8:  " ₪    90.00 ",   # ביטוח בריאות (unchanged)
    9:  " ₪   600.00 ",   # ביטוח צד ג' (unchanged)
    10: " ₪ 12,638.26 ",  # סה"כ עלות שכר
    11: " ₪    52.66 ",   # עלות עובד שעתית (12,638.26/240)
}

# Bottom table — actual clients from 03/2026
# Columns (right→left in slide): רווח/לקוח, רווח/עובד, עלות/עובד, הכנסות/עובד, שעות, הסעות, מגורים, תשלום/ש, כמות, לקוח
# Cost base = 52.66 × 240 = 12,638.40 ≈ 12,638
# Cost for ולפמן = (52.66 + 2.8 + 1.5) × 240 = 56.96 × 240 = 13,670
CLIENTS = [
    # (לקוח, כמות, תשלום/ש, מגורים, הסעות)
    ("מ.אילון אביב נכסים",  90, 65, "-", "-"),
    ("ולפמן תעשיות בע\"מ",  47, 70, "2.8", "1.5"),
    ("נח רפפורט 1990 בע\"מ", 20, 63, "-", "-"),
    ("א.ש עבודות בידוד",    15, 75, "-", "-"),
    ("קבוצת טלאור כראדי",    9, 75, "-", "-"),
    ("סלמאן והבה ובניו",     5, 75, "-", "-"),
]

BASE_COST_PER_HOUR = 52.66   # ₪ employer cost per hour
STANDARD_HOURS     = 240


def calc_row(client, count, rate, housing_str, transport_str):
    """Return (הכנסות/עובד, עלות/עובד, רווח/עובד, רווח/לקוח)."""
    housing   = float(housing_str)   if housing_str   != "-" else 0.0
    transport = float(transport_str) if transport_str != "-" else 0.0
    revenue   = rate * STANDARD_HOURS
    cost      = (BASE_COST_PER_HOUR + housing + transport) * STANDARD_HOURS
    profit    = revenue - cost
    total     = profit * count
    return int(round(revenue)), int(round(cost)), int(round(profit)), int(round(total))


def fmt(n):
    """Format integer with thousands separator."""
    return f"{n:,}"


def set_cell_text(cell, new_text: str):
    """Replace all runs in cell's first paragraph with a single run containing new_text.
    Preserves the formatting of the first run; removes extra runs."""
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if not para.runs:
        return
    # Set first run text
    para.runs[0].text = new_text
    # Remove any additional runs from the XML (keep only the first <a:r>)
    p_elem = para._p
    runs_xml = p_elem.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
    )
    for extra_run in runs_xml[1:]:
        p_elem.remove(extra_run)


def main():
    prs   = Presentation(SRC)
    slide = prs.slides[17]   # 0-based index → slide 18

    shapes = {s.name: s for s in slide.shapes}
    t_top = shapes["Table 8"].table
    t_bot = shapes["Table 10"].table

    # ── Update TOP table row 2 (עובד ממוצע) ─────────────────────────────────
    for col, text in TOP_ROW.items():
        set_cell_text(t_top.cell(2, col), text)
    print("✓ Updated top table (average employee costs)")

    # ── Update BOTTOM table rows 1-6 ────────────────────────────────────────
    # Columns (index 0-9): 0=רווח/לקוח, 1=רווח/עובד, 2=עלות/עובד, 3=הכנסות/עובד,
    #                      4=שעות, 5=הסעות, 6=מגורים, 7=תשלום/ש, 8=כמות, 9=לקוח
    for row_idx, (client, count, rate, housing, transport) in enumerate(CLIENTS, start=1):
        revenue, cost, profit, total = calc_row(client, count, rate, housing, transport)
        updates = {
            0: fmt(total),
            1: fmt(profit),
            2: fmt(cost),
            3: fmt(revenue),
            4: str(STANDARD_HOURS),
            5: transport,
            6: housing,
            7: str(rate),
            8: str(count),
            9: client,
        }
        for col, text in updates.items():
            set_cell_text(t_bot.cell(row_idx, col), text)
        print(f"  Row {row_idx}: {client} ({count} עובדים) — רווח חודשי: ₪{fmt(total)}")

    print()
    prs.save(DEST)
    print(f"✓ Saved: {DEST}")


if __name__ == "__main__":
    main()
